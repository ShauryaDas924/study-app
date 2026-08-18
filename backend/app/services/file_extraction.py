
import base64
import asyncio
from openai import OpenAI
from PIL import Image
import pymupdf as fitz
from pptx import Presentation
from io import BytesIO
from dotenv import load_dotenv
import os
import re
import zipfile

load_dotenv()


def _bounded_int_env(name: str, default: int, minimum: int, maximum: int) -> int:
    raw = os.getenv(name, str(default))
    try:
        value = int(raw)
    except ValueError as exc:
        raise RuntimeError(f"{name} must be an integer.") from exc
    if not minimum <= value <= maximum:
        raise RuntimeError(f"{name} must be between {minimum} and {maximum}.")
    return value


USE_VISION = True
MAX_PDF_PAGES = _bounded_int_env("MAX_PDF_PAGES", 100, 1, 500)
MAX_VISION_OCR_PAGES = _bounded_int_env("MAX_VISION_OCR_PAGES", 10, 0, 50)
MAX_IMAGE_PIXELS = _bounded_int_env("MAX_IMAGE_PIXELS", 25_000_000, 1, 100_000_000)
MAX_ARCHIVE_UNCOMPRESSED_BYTES = _bounded_int_env(
    "MAX_ARCHIVE_UNCOMPRESSED_BYTES", 100 * 1024 * 1024, 1, 500 * 1024 * 1024
)
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


def validate_image_bytes(file_bytes: bytes) -> str:
    try:
        with Image.open(BytesIO(file_bytes)) as image:
            width, height = image.size
            if width <= 0 or height <= 0 or width * height > MAX_IMAGE_PIXELS:
                raise ValueError("Image dimensions exceed the configured safety limit.")
            if image.format not in {"PNG", "JPEG"}:
                raise ValueError("Only PNG and JPEG images are supported.")
            mime_type = "image/jpeg" if image.format == "JPEG" else "image/png"
            image.verify()
            return mime_type
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("The uploaded image is invalid or corrupted.") from exc


def validate_pdf_document(doc) -> None:
    if len(doc) == 0:
        raise ValueError("The uploaded PDF has no pages.")
    if len(doc) > MAX_PDF_PAGES:
        raise ValueError(f"PDF exceeds the {MAX_PDF_PAGES}-page limit.")


def validate_pdf_bytes(file_bytes: bytes) -> None:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise ValueError("The uploaded PDF is invalid or corrupted.") from exc
    try:
        validate_pdf_document(doc)
    finally:
        doc.close()


def validate_pdf_page_render(page, dpi: int = 300) -> None:
    width = max(0.0, float(page.rect.width)) * dpi / 72
    height = max(0.0, float(page.rect.height)) * dpi / 72
    if width * height > MAX_IMAGE_PIXELS:
        raise ValueError("A PDF page exceeds the configured render-pixel limit.")

# ------------------------
# VISION LLM
# ------------------------
async def extract_from_image(file_bytes: bytes) -> str:
    mime_type = validate_image_bytes(file_bytes)
    b64 = base64.b64encode(file_bytes).decode()

    resp = await asyncio.to_thread(
        client.chat.completions.create,
        model="gpt-4.1",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all notes clearly. Preserve math."},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime_type};base64,{b64}"},
                    },
                ],
            }
        ],
    )

    return resp.choices[0].message.content


# ------------------------
# PDF
# ------------------------
def extract_from_pdf(file_bytes: bytes) -> str:
    text = ""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        validate_pdf_document(doc)
        for page in doc:
            text += page.get_text()
        return text
    finally:
        doc.close()


# ------------------------
# PDF (MATH MODE OCR)
# ------------------------
async def extract_from_pdf_math(file_bytes: bytes) -> str:

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    ocr_pages = 0
    try:
        validate_pdf_document(doc)
        for page in doc:
            page_text = page.get_text()

            if (len(page_text) < 200 or "  " in page_text) and ocr_pages < MAX_VISION_OCR_PAGES:
                validate_pdf_page_render(page)
                pix = page.get_pixmap(dpi=300)
                page_text = await extract_from_image(pix.tobytes("png"))
                ocr_pages += 1

            text += str(page_text or "") + "\n"
        return text
    finally:
        doc.close()

# ------------------------
# PPT
# ------------------------
def extract_from_ppt(file_bytes: bytes) -> str:
    try:
        with zipfile.ZipFile(BytesIO(file_bytes)) as archive:
            expanded_size = sum(item.file_size for item in archive.infolist())
            if expanded_size > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                raise ValueError("Presentation expands beyond the configured safety limit.")
    except ValueError:
        raise
    except zipfile.BadZipFile as exc:
        raise ValueError("The uploaded presentation is invalid or corrupted.") from exc

    prs = Presentation(BytesIO(file_bytes))
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def extract_from_plain_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


# ------------------------
# MAIN ROUTER FUNCTION
# ------------------------
async def extract_text(
    filename: str,
    file_bytes: bytes,
    math_mode: bool = False,
    allow_vision_ocr: bool = True,
) -> str:
    filename = filename.lower()

    if filename.endswith((".png", ".jpg", ".jpeg")):
        if USE_VISION and allow_vision_ocr:
            return await extract_from_image(file_bytes)
        raise ValueError("Image extraction requires vision OCR to be enabled.")

    if filename.endswith(".pdf"):

        if math_mode and allow_vision_ocr:
            return await extract_from_pdf_math(file_bytes)

        return extract_from_pdf(file_bytes)

    if filename.endswith(".pptx"):
        return extract_from_ppt(file_bytes)

    if filename.endswith((".txt", ".md")):
        return extract_from_plain_text(file_bytes)

    raise ValueError("Unsupported file type.")


async def extract_text_with_source(
    filename: str,
    file_bytes: bytes,
    math_mode: bool = False,
    allow_vision_ocr: bool = True,
) -> dict:
    """
    Extract text while keeping lightweight source metadata for exam-prep evidence.
    Existing callers should keep using extract_text; this helper is additive.
    """
    original_filename = filename or "uploaded_material"
    lower = original_filename.lower()

    if lower.endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            validate_pdf_document(doc)
            page_count = len(doc)
            pages = []
            full_text_parts = []
            ocr_skipped_pages = []
            ocr_pages = 0

            for index, page in enumerate(doc, start=1):
                page_text = page.get_text()

                needs_ocr = math_mode and (len(page_text or "") < 200 or "  " in (page_text or ""))
                if needs_ocr and allow_vision_ocr and ocr_pages < MAX_VISION_OCR_PAGES:
                    validate_pdf_page_render(page)
                    pix = page.get_pixmap(dpi=300)
                    page_text = await extract_from_image(pix.tobytes("png"))
                    ocr_pages += 1
                elif needs_ocr:
                    ocr_skipped_pages.append(index)

                page_text = str(page_text or "").strip()
                if page_text:
                    pages.append(
                        {
                            "page": index,
                            "start_char": sum(len(part) + 1 for part in full_text_parts),
                            "text": page_text,
                        }
                    )
                    full_text_parts.append(f"[Page {index}]\n{page_text}")

            return {
                "text": "\n\n".join(full_text_parts).strip(),
                "pages": pages,
                "source_ref": {
                    "filename": original_filename,
                    "page_count": page_count,
                    "extraction_mode": "pdf_pages",
                    "vision_ocr_allowed": allow_vision_ocr,
                    "ocr_skipped_pages": ocr_skipped_pages,
                },
            }
        finally:
            doc.close()

    text = await extract_text(
        original_filename,
        file_bytes,
        math_mode=math_mode,
        allow_vision_ocr=allow_vision_ocr,
    )
    return {
        "text": text,
        "pages": [],
        "source_ref": {
            "filename": original_filename,
            "extraction_mode": "single_text_block",
            "vision_ocr_allowed": allow_vision_ocr,
        },
    }


# ------------------------
# HOMEWORK QUESTION SPLITTER
# ------------------------



def split_homework_questions(text: str):

    # Splits on patterns like:
    # 1)
    # 1.
    # Question 1
    # Q1

    pattern = r"(?:^|\n)(?:Question\s*\d+|\d+\)|\d+\.)"

    parts = re.split(pattern, text)

    questions = []

    for p in parts:
        p = p.strip()

        if len(p) > 20:
            questions.append(p)

    return questions
